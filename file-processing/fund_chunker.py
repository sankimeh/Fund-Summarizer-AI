import os
import re
import hashlib
from typing import List, Dict, Any

import fitz  # PyMuPDF
import pdfplumber
import pytesseract
from PIL import Image
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

# ----------- Section Header Heuristic -----------
SECTION_HEADER_REGEX = re.compile(
    r"^(?:\d+(?:\.\d+)*\.?\s+)?[A-Z][\w\s,&/()-]{3,}$", re.MULTILINE
)

def extract_section_title(text: str) -> str:
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    for line in lines:
        if len(line.split()) <= 12 and line[0].isupper():
            return line
    return "Unknown Section"

def find_section_boundaries(text: str) -> List[Dict[str, str]]:
    sections = []
    matches = list(SECTION_HEADER_REGEX.finditer(text))
    for idx, match in enumerate(matches):
        start = match.start()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
        section_text = text[start:end].strip()
        header = match.group().strip().rstrip(":")
        sections.append({"title": header, "text": section_text})
    return sections

# ----------- Topic Tagging Logic -----------
TOPIC_TAGS = {
    "fees": ["management fee", "expense", "ter", "incentive", "admin"],
    "liquidity": ["redemption", "liquidity", "lock-up", "withdrawal", "notice period"],
    "risk": ["risk", "volatility", "default", "loss", "cyber"],
    "performance": ["irr", "return", "yield", "benchmark"],
    "structure": ["lp", "llc", "trust", "legal structure"],
    "strategy": ["investment strategy", "asset allocation", "portfolio"],
    "redemption_terms": ["redemption", "notice", "lock-in", "exit window"]
}

def get_financial_tags(text: str) -> List[str]:
    text = text.lower()
    tags = []
    for tag, keywords in TOPIC_TAGS.items():
        if any(kw in text for kw in keywords):
            tags.append(tag)
    return tags or ["general"]

def flatten_table(table: List[List[Any]]) -> str:
    return "\n".join(
        [" | ".join(str(cell or "").strip() for cell in row) for row in table if any(cell and str(cell).strip() for cell in row)]
    )

def chunk_large_table(table, max_lines=20):
    for i in range(0, len(table), max_lines):
        yield table[i:i + max_lines]

# ----------- PDF and PPT Extraction -----------
def extract_pdf_text_by_page(file_path: str, image_output_dir: str = "pdf_images") -> List[Dict[str, Any]]:
    os.makedirs(image_output_dir, exist_ok=True)
    pages = []

    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            page_data = {"page_number": i, "text": "", "tables": [], "is_ocr": False}
            text = page.extract_text()
            if not text or len(text.splitlines()) < 3:
                text = ocr_pdf_page(file_path, i - 1, image_output_dir)
                page_data["is_ocr"] = True
            page_data["text"] = text.strip() if text else ""
            tables = page.extract_tables()
            for table in tables:
                if table:
                    page_data["tables"].append(table)
            save_pdf_page_as_image(file_path, i - 1, os.path.join(image_output_dir, f"{os.path.basename(file_path)}_page_{i}.png"))
            pages.append(page_data)

    pages = detect_and_merge_multi_page_tables(pages)
    return pages

def detect_and_merge_multi_page_tables(pages: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    for i in range(len(pages) - 1):
        current_tables = pages[i]["tables"]
        next_tables = pages[i + 1]["tables"]
        if current_tables and next_tables:
            last_table = current_tables[-1]
            first_next_table = next_tables[0]
            if is_probable_continuation(last_table, first_next_table):
                merged = last_table + first_next_table
                pages[i]["tables"][-1] = merged
                del pages[i + 1]["tables"][0]
    return pages

def is_probable_continuation(last_table: List[List[str]], next_table: List[List[str]]) -> bool:
    if not last_table or not next_table:
        return False
    header_keywords = {"date", "amount", "description", "investment"}
    row = " ".join(str(cell or "") for cell in next_table[0]).lower()
    if any(kw in row for kw in header_keywords):
        return False
    return len(last_table) > 2 and len(next_table) > 2

def ocr_pdf_page(file_path: str, page_index: int, image_output_dir: str) -> str:
    with fitz.open(file_path) as doc:
        page = doc.load_page(page_index)
        pix = page.get_pixmap(dpi=300)
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        img_path = os.path.join(image_output_dir, f"{os.path.basename(file_path)}_page_{page_index+1}_ocr.png")
        img.save(img_path)
    return pytesseract.image_to_string(img)

def save_pdf_page_as_image(file_path: str, page_index: int, output_path: str):
    with fitz.open(file_path) as doc:
        page = doc.load_page(page_index)
        pix = page.get_pixmap(dpi=150)
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        img.save(output_path)

def extract_ppt_text_by_slide(file_path: str) -> List[Dict[str, Any]]:
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_data = {"page_number": i, "text": "", "tables": [], "is_ocr": False}
        text_lines = [f"Slide {i}:"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_lines.append(shape.text.strip())
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_text = [cell.text_frame.text if cell.text_frame else "" for cell in row.cells]
                    table_data.append(row_text)
                slide_data["tables"].append(table_data)
        slide_data["text"] = "\n".join(text_lines).strip()
        slides.append(slide_data)
    return slides

# ----------- Chunking Logic -----------
def chunk_pages(
    pages: List[Dict[str, Any]],
    fund_name: str,
    source_file: str,
    chunk_size: int = 1000,
    overlap: int = 150,
    mode: str = "auto"
) -> List[Dict[str, Any]]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ".", " "]
    )

    chunks = []
    for page in pages:
        raw_text = page["text"]
        page_number = page["page_number"]
        has_table = bool(page["tables"])

        # Append tables text as chunks too
        if has_table:
            for table in page["tables"]:
                # chunk_large_table yields smaller slices of table rows
                for table_chunk in chunk_large_table(table):
                    table_text = flatten_table(table_chunk)
                    raw_text += f"\n\n[TABLE DATA]\n{table_text}"

        section_chunks = []
        if mode == "character":
            text_chunks = splitter.split_text(raw_text)
            for chunk in text_chunks:
                section_chunks.append({
                    "title": extract_section_title(chunk),
                    "text": chunk
                })
        else:
            found_sections = find_section_boundaries(raw_text)
            if not found_sections:
                text_chunks = splitter.split_text(raw_text)
                for chunk in text_chunks:
                    section_chunks.append({
                        "title": extract_section_title(chunk),
                        "text": chunk
                    })
            else:
                section_chunks = found_sections

        for i, sec in enumerate(section_chunks):
            section_text = sec["text"]
            metadata = {
                "fund_name": fund_name,
                "section_title": sec["title"],
                "page_number": page_number,
                "source_file": source_file,
                "prev_section": section_chunks[i - 1]["title"] if i > 0 else "",
                "next_section": section_chunks[i + 1]["title"] if i < len(section_chunks) - 1 else "",
                "chunk_index": i,
                "financial_tags": ", ".join(get_financial_tags(section_text)),  # <-- HERE
                "is_ocr": page.get("is_ocr", False),
                "has_table": has_table,
                "source_checksum": hashlib.md5(section_text.encode("utf-8")).hexdigest()[:10]
            }
            chunks.append({
                "text": section_text,
                "metadata": metadata
            })

    return chunks

# ----------- Controller Function -----------
def process_files(
    fund_name: str,
    file_paths: List[str],
    image_output_dir: str = "pdf_images",
    chunk_size: int = 1000,
    overlap: int = 150,
    chunk_mode: str = "auto"
) -> List[Dict[str, Any]]:
    all_chunks = []
    for file_path in file_paths:
        ext = os.path.splitext(file_path)[1].lower()
        source_file = os.path.basename(file_path)

        if ext == ".pdf":
            pages = extract_pdf_text_by_page(file_path, image_output_dir)
        elif ext in {".ppt", ".pptx"}:
            pages = extract_ppt_text_by_slide(file_path)
        else:
            print(f"[Skipped] Unsupported file format: {file_path}")
            continue

        chunks = chunk_pages(
            pages,
            fund_name=fund_name,
            source_file=source_file,
            chunk_size=chunk_size,
            overlap=overlap,
            mode=chunk_mode
        )
        all_chunks.extend(chunks)

    return all_chunks
